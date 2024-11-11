import tempfile
import io
import zipfile
import os
import re
import shutil
from pptx import Presentation
from openpyxl import load_workbook, Workbook

def analyze_excel_markers(workbook):
    """Analyze the Excel workbook to get all markers (sheet name, start, and end) and create a mapping."""
    mapping = {}
    try:
        for sheet in workbook.worksheets:
            start_row = None
            end_row = None
            sheet_name = sheet.title

            for row_idx, row in enumerate(sheet.iter_rows(min_col=1, max_col=1, values_only=True), start=1):
                cell_value = row[0]
                if isinstance(cell_value, str) and cell_value.startswith("pptstart:"):
                    marker_name = cell_value.split(":")[1]
                    start_row = row_idx
                elif isinstance(cell_value, str) and cell_value.startswith("pptend:"):
                    end_row = row_idx - 1

                if start_row and end_row:
                    mapping[marker_name] = {
                        "sheet_name": sheet_name,
                        "start_row": start_row,
                        "end_row": end_row
                    }
                    start_row = None
                    end_row = None

    except Exception as e:
        print(f"Failed to analyze Excel workbook: {e}")

    return mapping

def copy_cells_into_embedded_workbook(source_workbook, mapping, marker_name, embedded_sheet_name):
    """Create a new workbook and copy cells from the source Excel sheet based on the marker."""
    try:
        if marker_name not in mapping:
            print(f"Marker '{marker_name}' not found in the Excel mapping. Skipping chart.")
            return None  # Return None if marker not found

        mapping_info = mapping[marker_name]
        source_sheet_name = mapping_info["sheet_name"]
        source_start_row = mapping_info["start_row"]
        source_end_row = mapping_info["end_row"]

        source_worksheet = source_workbook[source_sheet_name]
        
        # Create a new workbook object
        new_workbook = Workbook()
        new_worksheet = new_workbook.create_sheet(title=embedded_sheet_name)
        new_workbook.remove(new_workbook.active)

        for source_row_idx, row in enumerate(source_worksheet.iter_rows(min_row=source_start_row, max_row=source_end_row, values_only=True), start=1):
            for col_idx, cell_value in enumerate(row, start=1):
                new_worksheet.cell(row=source_row_idx, column=col_idx).value = cell_value

        return new_workbook  # Return the new workbook

    except Exception as e:
        print(f"Failed to copy data for marker '{marker_name}': {e}")
        return None

def process_embedded_workbook(embedded_file, embedded_content, workbook, mapping):
    """Process each embedded Excel workbook from in-memory ZIP."""
    try:
        # embedded_workbook = load_workbook(io.BytesIO(embedded_content), read_only=True)
        embedded_workbook = load_workbook(io.BytesIO(embedded_content), read_only=True, data_only=True)
        # embedded_workbook = load_workbook(io.BytesIO(embedded_content))
        sheet_embedded = embedded_workbook.active
        embedded_sheet_name = sheet_embedded.title

        marker = sheet_embedded['A1'].value
        if marker is None or not isinstance(marker, str) or not marker.startswith("pptstart:"):
            print(f"No valid marker found in A1 of the embedded workbook '{embedded_file}'. Skipping.")
            return None

        marker_name = marker.split(":")[1]

        new_workbook = copy_cells_into_embedded_workbook(workbook, mapping, marker_name, embedded_sheet_name)

        if new_workbook is None:
            return None

        # Save the new workbook to memory
        new_workbook_io = io.BytesIO()
        new_workbook.save(new_workbook_io)
        new_workbook_io.seek(0)  # Move cursor back to the beginning of the stream

        return embedded_file, new_workbook_io.read()  # Return new workbook content

    except Exception as e:
        print(f"Failed to modify '{embedded_file}': {e}")
        return None

def modify_embedded_excel_in_pptx(config):
    """Modify embedded Excel files in the PowerPoint presentation based on the provided mapping."""
    ppt_file_path = config['ppt_file_path']
    workbook = config['workbook']
    mapping = config['mapping']
    use_filesystem = config.get('use_filesystem', False)
    temp_dir = None
    try:
        if use_filesystem:
            # Use unique temporary directory for filesystem processing
            temp_dir = tempfile.mkdtemp()

            with zipfile.ZipFile(ppt_file_path, 'r') as ppt_zip:
                ppt_zip.extractall(temp_dir)

                for item in ppt_zip.infolist():
                    if item.filename.startswith('ppt/embeddings/') and item.filename.endswith('.xlsx'):
                        embedded_path = os.path.join(temp_dir, item.filename)
                        with open(embedded_path, 'rb') as f:
                            embedded_content = f.read()
                        result = process_embedded_workbook(item.filename, embedded_content, workbook, mapping)
                        if result:
                            embedded_filename, new_workbook_content = result
                            with open(os.path.join(temp_dir, embedded_filename), 'wb') as f:
                                f.write(new_workbook_content)

                # Replace placeholders in slide XML files
                slide_dir = os.path.join(temp_dir, 'ppt/slides')
                if os.path.exists(slide_dir):
                    for slide_file in os.listdir(slide_dir):
                        slide_path = os.path.join(slide_dir, slide_file)
                        if slide_file.endswith('.xml'):
                            with open(slide_path, 'r', encoding='utf-8') as file:
                                slide_content = file.read()

                            updated_slide_content = replace_placeholders_in_slide_content(slide_content, workbook)

                            with open(slide_path, 'w', encoding='utf-8') as file:
                                file.write(updated_slide_content)

            updated_ppt_io = io.BytesIO()
            with zipfile.ZipFile(updated_ppt_io, 'w') as updated_ppt_zip:
                for root, dirs, files in os.walk(temp_dir):
                    for file in files:
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, temp_dir)
                        updated_ppt_zip.write(file_path, arcname)

            updated_ppt_io.seek(0)

            presentation = Presentation(updated_ppt_io)

        else:
            # Use in-memory zipping/unzipping
            with zipfile.ZipFile(ppt_file_path, 'r') as ppt_zip:
                updated_ppt_io = io.BytesIO()
                updated_ppt_zip = zipfile.ZipFile(updated_ppt_io, 'w')

                for item in ppt_zip.infolist():
                    if item.filename.startswith('ppt/embeddings/') and item.filename.endswith('.xlsx'):
                        embedded_content = ppt_zip.read(item.filename)
                        result = process_embedded_workbook(item.filename, embedded_content, workbook, mapping)

                        if result:
                            embedded_filename, new_workbook_content = result
                            updated_ppt_zip.writestr(embedded_filename, new_workbook_content)
                        else:
                            updated_ppt_zip.writestr(item.filename, embedded_content)
                    elif item.filename.startswith('ppt/slides') and item.filename.endswith('.xml'):
                        slide_content = ppt_zip.read(item.filename).decode('utf-8')
                        updated_slide_content = replace_placeholders_in_slide_content(slide_content, workbook)
                        updated_ppt_zip.writestr(item.filename, updated_slide_content.encode('utf-8'))
                    else:
                        updated_ppt_zip.writestr(item.filename, ppt_zip.read(item.filename))

                updated_ppt_zip.close()
                updated_ppt_io.seek(0)

                presentation = Presentation(updated_ppt_io)

    except Exception as e:
        print(f"An error occurred: {e}")
        raise
    finally:
        # Clean up the temporary directory if it was used
        if temp_dir and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)

    return presentation

def get_excel_value(sheet, cell_ref):
    """
    Fetch the value from the specified Excel sheet and cell reference.
    """
    try:
        cell = sheet[cell_ref]
        return format_value(cell.value, cell.number_format)
    except Exception as e:
        print(f"Error occurred: {str(e)}")
        return ""

def format_value(value, number_format):
    """
    Format the value according to the Excel cell's number format.
    - Retain Excel's significant digits and percentage formatting.
    - Always format numbers over 1 million with "M".
    """
    if value is None:
        return " "  # Return a space instead of "N/A"

    # Handle numbers over 1 million, format them with "M"
    if isinstance(value, (int, float)) and value >= 1_000_000:
        return f"{value / 1_000_000:.1f}M"

    # Handle percentage formats based on Excel formatting
    if "0%" in number_format or "percent" in number_format.lower():
        decimals = number_format.split(".")[1].count("0") if "." in number_format else 0
        return f"{value * 100:.{decimals}f}%"

    # Handle significant digits based on Excel's formatting
    if isinstance(value, (int, float)):
        if number_format == "General":
            return str(value)
        elif "0." in number_format:
            decimals = number_format.split(".")[1].count("0")  # Count decimal places
            return f"{value:.{decimals}f}"
        elif number_format.startswith("#,##0"):
            return f"{value:,.0f}"
        elif number_format == "0":
            return f"{int(value)}"

    return str(value)

def replace_placeholders_in_slide_content(slide_content, workbook):
    """Replace placeholders in the slide content using a regex."""
    sheets = {sheet.title: sheet for sheet in workbook}
    placeholder_pattern = re.compile(r'\[\[([A-Za-z0-9_ ]+?)!(\w+)\]\]')

    def replace_placeholder(match):
        sheet_name, cell_ref = match.groups()
        if sheet_name in sheets:
            return get_excel_value(sheets[sheet_name], cell_ref)
        return "Not found"

    updated_content = placeholder_pattern.sub(replace_placeholder, slide_content)
    return updated_content

